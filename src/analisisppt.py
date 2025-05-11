from src.pineconeupload import PineconeUploader
from pptx import Presentation
import os
from langchain.text_splitter import RecursiveCharacterTextSplitter

class PPTXProcessor:
    def __init__(self, pptx_dir):
        """Inicializa la clase con el directorio donde se encuentran los archivos .pptx."""
        self.pptx_dir = pptx_dir
        self.pcuploader = PineconeUploader()

    def list_pptx(self):
        """Lista todos los archivos .pptx en el directorio especificado."""
        archivos_pptx = [archivo for archivo in os.listdir(self.pptx_dir) if archivo.lower().endswith('.pptx')]
        return archivos_pptx

    def extract_text(self, fichero, charleados):
        """Extrae texto de las diapositivas de un archivo .pptx y organiza el contenido."""
        presentation = Presentation(os.path.join(self.pptx_dir, fichero))
        textos = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if len(text) > 2:  # Filtro para omitir textos breves o vacíos
                        
                        divided_texts = self.divide_text(text)
                
                        for text in divided_texts:
                            #self.pcuploader.upload_text("Wiki: "+pagina  + " \n\n"+text, rol="SYSTEM") 

                            #print(fichero)
                            
                            if fichero in charleados: #TODO revisar esto
                                
                                texto_formateado = " ".join(charleados[fichero])+ f" {fichero}: Diapositiva {{{slide_number}}}\n{text}" 
                            
                            else:
                                texto_formateado = f"{fichero}: Diapositiva {{{slide_number}}}\n{text}"
                            textos.append(texto_formateado)
        
        
        
     
        #return res #condensarlos todos en 1     
        return textos 

    
    def divide_text(self, text):
        #texts = []
        
        texts = RecursiveCharacterTextSplitter(chunk_size=8180, chunk_overlap=100).split_text(text) #https://medium.com/@vndee.huynh/build-your-own-rag-and-run-it-locally-langchain-ollama-streamlit-181d42805895
        # sacado de https://huggingface.co/jinaai/jina-embeddings-v2-base-es 
        # un poco menos para que entre lo la parte de get_embedding
        #texts = filter_complex_metadata(texts)
        
        return texts
    
    
    def analyze_and_upload(self, charleados,rol="SYSTEM"):
        """Procesa todos los archivos .pptx en el directorio y sube el contenido extraído a Pinecone."""
        ficheros = self.list_pptx()
        res = ""
        for fichero in ficheros:
            textos = self.extract_text(fichero, charleados)
            if textos:  # Solo subir si hay textos válidos
                """""
                for texto in textos:
                    res=res + " " + texto
                """
                #self.pcuploader.upload_text(res, "SYSTEM")
                self.pcuploader.bulk_upload(textos, rol=rol)
                print(f"Contenido de {fichero} subido correctamente.")
            else:
                print(f"No se encontró contenido válido en {fichero}.")

